from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

mainObj = dict()
# slideObj = dict()
# subText = list()
prs = Presentation(".pptx")

for index, slide in enumerate(prs.slides):
    for indexSlide, shape in enumerate(slide.shapes):
        if shape.is_placeholder:
            phf = shape.placeholder_format
            # print('%d, %s' % (phf.idx, phf.type))
            if phf.type == PP_PLACEHOLDER.TITLE or phf.type == PP_PLACEHOLDER.CENTER_TITLE or phf.type == PP_PLACEHOLDER.VERTICAL_TITLE:
                print phf.type, shape.text_frame.text
                mainObj[index] = shape.text_frame.text
                # print('%d %s' % (shape.placeholder_format.idx, shape.name))
                # print shape.shape_type
                # # print shape.shape_type
                # if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                #     print shape
                # except AttributeError:
                #     pass

                # slideObj['body'] = subText[0]
                # mainObj[index] = slideObj
                # del subText[:]

print mainObj

# can be also done for subtitles
