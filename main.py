from pptx import Presentation

mainObj = dict()
slideObj = dict()
subText = list()
prs = Presentation("microprocessor.pptx")


for index, slide in enumerate(prs.slides):
    for indexSlide, shape in enumerate(slide.shapes):
        try:
            if indexSlide == 0:
                slideObj['title'] = shape.text
            else:
                subText.append(shape.text)
        except AttributeError:
            pass

    slideObj['body'] = subText[0]
    mainObj[index] = slideObj
    del subText[:]

print mainObj
